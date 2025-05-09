"""
File connector for retrieving content from various file types.
"""

import logging
import os
import glob
import json
import re
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Generator, Union, Set

from ..connectors.base import BaseConnector, Document

logger = logging.getLogger(__name__)


class FileConnector(BaseConnector):
    """Connector for retrieving content from files of various types."""
    
    SUPPORTED_EXTENSIONS = {
        "text": [".txt", ".csv", ".tsv", ".log"],
        "markdown": [".md", ".markdown"],
        "pdf": [".pdf"],
        "docx": [".docx", ".doc"],
        "pptx": [".pptx", ".ppt"],
        "xlsx": [".xlsx", ".xls"],
        "json": [".json", ".jsonl"],
        "xml": [".xml"],
        "html": [".html", ".htm"],
        "email": [".eml", ".msg"],
        "code": [".py", ".js", ".java", ".cpp", ".c", ".cs", ".php", 
                ".rb", ".go", ".rs", ".ts", ".sh", ".ps1", ".bat"]
    }
    
    def __init__(self, config: Dict[str, Any]):
        """
        Initialize the file connector.
        
        Args:
            config: Dictionary with the following keys:
                - base_dirs: List of directories to scan for files
                - include_patterns: List of glob patterns to include
                - exclude_patterns: List of glob patterns to exclude
                - recursive: Whether to scan directories recursively
                - follow_symlinks: Whether to follow symbolic links
                - max_file_size_mb: Maximum file size in MB
                - max_files: Maximum number of files to process
                - extensions: List of file extensions to process
                - extract_metadata: Whether to extract metadata from files
        """
        super().__init__(config)
        self.base_dirs = [Path(p) for p in config.get("base_dirs", [])]
        self.include_patterns = config.get("include_patterns", ["*.*"])
        self.exclude_patterns = config.get("exclude_patterns", [".*", "*.pyc", "*.tmp"])
        self.recursive = config.get("recursive", True)
        self.follow_symlinks = config.get("follow_symlinks", False)
        self.max_file_size = config.get("max_file_size_mb", 10) * 1024 * 1024  # Convert to bytes
        self.max_files = config.get("max_files", 1000)
        self.extensions = config.get("extensions", list(self.SUPPORTED_EXTENSIONS.keys()))
        self.extract_metadata = config.get("extract_metadata", True)
        
        # Validate configuration
        if not self.base_dirs:
            raise ValueError("File connector requires at least one base directory")
        
        # Initialize supported file types
        self._parsers = {}
        self._load_parsers()
    
    def _load_parsers(self):
        """Initialize file parsers for each supported file type."""
        # Standard text/markdown files
        self._parsers["text"] = self._extract_text_content
        self._parsers["markdown"] = self._extract_text_content
        
        # JSON files
        self._parsers["json"] = self._extract_json_content
        
        # PDF files
        try:
            import pypdf
            self._parsers["pdf"] = self._extract_pdf_content
        except ImportError:
            logger.warning("pypdf not installed, PDF extraction disabled")
        
        # DOCX files
        try:
            import docx2txt
            self._parsers["docx"] = self._extract_docx_content
        except ImportError:
            logger.warning("docx2txt not installed, DOCX extraction disabled")
        
        # PPTX files
        try:
            from pptx import Presentation
            self._parsers["pptx"] = self._extract_pptx_content
        except ImportError:
            logger.warning("python-pptx not installed, PPTX extraction disabled")
        
        # XLSX files
        try:
            import pandas as pd
            self._parsers["xlsx"] = self._extract_xlsx_content
        except ImportError:
            logger.warning("pandas not installed, XLSX extraction disabled")
        
        # HTML files
        try:
            from bs4 import BeautifulSoup
            import html2text
            self._parsers["html"] = self._extract_html_content
        except ImportError:
            logger.warning("beautifulsoup4/html2text not installed, HTML extraction disabled")
        
        # XML files
        try:
            import xmltodict
            self._parsers["xml"] = self._extract_xml_content
        except ImportError:
            logger.warning("xmltodict not installed, XML extraction disabled")
        
        # Email files
        try:
            import email
            self._parsers["email"] = self._extract_email_content
        except ImportError:
            logger.warning("email package not available, EML extraction disabled")
    
    async def connect(self) -> bool:
        """
        Verify access to the base directories.
        
        Returns:
            True if all directories are accessible, False otherwise
        """
        for base_dir in self.base_dirs:
            if not os.path.isdir(base_dir):
                logger.error(f"Base directory not found: {base_dir}")
                return False
            if not os.access(base_dir, os.R_OK):
                logger.error(f"No read permission for directory: {base_dir}")
                return False
        return True
    
    async def disconnect(self) -> None:
        """No need to disconnect for file operations."""
        pass
    
    def _get_all_files(self) -> List[Path]:
        """Get list of all files matching the patterns."""
        all_files = []
        
        for base_dir in self.base_dirs:
            # Skip if base directory doesn't exist
            if not os.path.isdir(base_dir):
                continue
            
            # Get the include extensions list
            include_extensions = []
            for ext_type in self.extensions:
                if ext_type in self.SUPPORTED_EXTENSIONS:
                    include_extensions.extend(self.SUPPORTED_EXTENSIONS[ext_type])
            
            # Process include patterns
            for pattern in self.include_patterns:
                # Adjust pattern with base directory
                full_pattern = os.path.join(base_dir, pattern)
                
                # Find files matching the pattern
                for file_path in glob.glob(full_pattern, recursive=self.recursive):
                    path = Path(file_path)
                    if path.is_file():
                        # Check if it's a supported extension
                        if include_extensions and path.suffix.lower() not in include_extensions:
                            continue
                        
                        # Skip excluded files
                        if any(path.match(exclude) for exclude in self.exclude_patterns):
                            continue
                        
                        # Skip symlinks if not following them
                        if not self.follow_symlinks and path.is_symlink():
                            continue
                        
                        # Skip files that are too large
                        if path.stat().st_size > self.max_file_size:
                            continue
                        
                        all_files.append(path)
            
            # Stop if we've found enough files
            if len(all_files) >= self.max_files:
                break
        
        return all_files[:self.max_files]
    
    def _get_file_type(self, file_path: Path) -> str:
        """Determine the file type from its extension."""
        suffix = file_path.suffix.lower()
        
        for file_type, extensions in self.SUPPORTED_EXTENSIONS.items():
            if suffix in extensions:
                return file_type
        
        # Default to treating unknown extensions as text
        return "text"
    
    def _extract_text_content(self, file_path: Path) -> str:
        """Extract content from plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error reading text file {file_path}: {e}")
            return ""
    
    def _extract_json_content(self, file_path: Path) -> str:
        """Extract content from JSON files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                # Check if it's a JSON Lines file
                if file_path.suffix.lower() == '.jsonl':
                    # Process as JSONL
                    lines = []
                    for i, line in enumerate(f):
                        if i >= 100:  # Limit to first 100 lines
                            lines.append("... (truncated)")
                            break
                        try:
                            data = json.loads(line)
                            lines.append(json.dumps(data, indent=2))
                        except json.JSONDecodeError:
                            lines.append(line)
                    return "\n".join(lines)
                else:
                    # Process as regular JSON
                    data = json.load(f)
                    return json.dumps(data, indent=2)
        except json.JSONDecodeError as e:
            logger.error(f"Error parsing JSON file {file_path}: {e}")
            # Fall back to reading as text
            return self._extract_text_content(file_path)
        except Exception as e:
            logger.error(f"Error reading JSON file {file_path}: {e}")
            return ""
    
    def _extract_pdf_content(self, file_path: Path) -> str:
        """Extract content from PDF files."""
        try:
            import pypdf
            with open(file_path, 'rb') as f:
                pdf = pypdf.PdfReader(f)
                text = []
                for page_num in range(len(pdf.pages)):
                    page = pdf.pages[page_num]
                    text.append(f"--- Page {page_num + 1} ---")
                    text.append(page.extract_text())
                return "\n".join(text)
        except Exception as e:
            logger.error(f"Error extracting content from PDF {file_path}: {e}")
            return ""
    
    def _extract_docx_content(self, file_path: Path) -> str:
        """Extract content from DOCX files."""
        try:
            import docx2txt
            text = docx2txt.process(file_path)
            return text
        except Exception as e:
            logger.error(f"Error extracting content from DOCX {file_path}: {e}")
            return ""
    
    def _extract_pptx_content(self, file_path: Path) -> str:
        """Extract content from PPTX files."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            
            for i, slide in enumerate(prs.slides):
                text.append(f"--- Slide {i + 1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text.append(shape.text)
            
            return "\n".join(text)
        except Exception as e:
            logger.error(f"Error extracting content from PPTX {file_path}: {e}")
            return ""
    
    def _extract_xlsx_content(self, file_path: Path) -> str:
        """Extract content from XLSX files."""
        try:
            import pandas as pd
            
            # Read all sheets
            sheets = pd.read_excel(file_path, sheet_name=None)
            text = []
            
            for sheet_name, df in sheets.items():
                text.append(f"--- Sheet: {sheet_name} ---")
                
                # Add column headers
                text.append("| " + " | ".join(str(col) for col in df.columns) + " |")
                text.append("| " + " | ".join("---" for _ in df.columns) + " |")
                
                # Add rows (limit to 100 rows)
                for i, row in df.head(100).iterrows():
                    text.append("| " + " | ".join(str(val) for val in row) + " |")
                
                if len(df) > 100:
                    text.append(f"... (showing 100 of {len(df)} rows)")
            
            return "\n".join(text)
        except Exception as e:
            logger.error(f"Error extracting content from XLSX {file_path}: {e}")
            return ""
    
    def _extract_html_content(self, file_path: Path) -> str:
        """Extract content from HTML files."""
        try:
            from bs4 import BeautifulSoup
            import html2text
            
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                html = f.read()
            
            # First try to extract with BeautifulSoup
            soup = BeautifulSoup(html, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.extract()
            
            # Extract title
            title = soup.title.text if soup.title else ""
            
            # Convert to markdown using html2text
            converter = html2text.HTML2Text()
            converter.ignore_links = False
            converter.ignore_images = False
            markdown = converter.handle(html)
            
            if title:
                return f"Title: {title}\n\n{markdown}"
            else:
                return markdown
        except Exception as e:
            logger.error(f"Error extracting content from HTML {file_path}: {e}")
            # Fall back to reading as text
            return self._extract_text_content(file_path)
    
    def _extract_xml_content(self, file_path: Path) -> str:
        """Extract content from XML files."""
        try:
            import xmltodict
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                xml = f.read()
            data = xmltodict.parse(xml)
            return json.dumps(data, indent=2)
        except Exception as e:
            logger.error(f"Error extracting content from XML {file_path}: {e}")
            # Fall back to reading as text
            return self._extract_text_content(file_path)
    
    def _extract_email_content(self, file_path: Path) -> str:
        """Extract content from email files."""
        try:
            import email
            from email.header import decode_header
            
            with open(file_path, 'rb') as f:
                msg = email.message_from_binary_file(f)
            
            # Helper function to decode email headers
            def decode_str(s):
                if s is None:
                    return ""
                decoded_parts = []
                for part, encoding in decode_header(s):
                    if isinstance(part, bytes):
                        try:
                            if encoding:
                                decoded_parts.append(part.decode(encoding))
                            else:
                                decoded_parts.append(part.decode())
                        except (UnicodeDecodeError, LookupError):
                            decoded_parts.append(part.decode('latin1', errors='replace'))
                    else:
                        decoded_parts.append(part)
                return "".join(decoded_parts)
            
            # Extract headers
            subject = decode_str(msg["Subject"])
            from_addr = decode_str(msg["From"])
            to_addr = decode_str(msg["To"])
            date = decode_str(msg["Date"])
            
            # Extract body
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    content_type = part.get_content_type()
                    content_disposition = str(part.get("Content-Disposition"))
                    
                    # Skip attachments
                    if "attachment" in content_disposition:
                        continue
                    
                    try:
                        # Get the body
                        payload = part.get_payload(decode=True)
                        if payload:
                            charset = part.get_content_charset()
                            if not charset:
                                charset = 'utf-8'
                                
                            try:
                                decoded_payload = payload.decode(charset)
                            except UnicodeDecodeError:
                                decoded_payload = payload.decode('latin1', errors='replace')
                            
                            if content_type == "text/plain":
                                body += decoded_payload
                            elif content_type == "text/html" and not body:
                                # Try to convert HTML to text
                                try:
                                    from bs4 import BeautifulSoup
                                    soup = BeautifulSoup(decoded_payload, 'html.parser')
                                    html_body = soup.get_text(separator="\n")
                                    if not body:  # Only use HTML if we don't have text yet
                                        body = html_body
                                except ImportError:
                                    # If BeautifulSoup isn't available, use a basic regex approach
                                    html_body = re.sub(r'<[^>]+>', ' ', decoded_payload)
                                    if not body:
                                        body = html_body
                    except Exception as e:
                        logger.error(f"Error extracting email part: {e}")
            else:
                # Not multipart - get the payload directly
                payload = msg.get_payload(decode=True)
                if payload:
                    charset = msg.get_content_charset()
                    if not charset:
                        charset = 'utf-8'
                        
                    try:
                        body = payload.decode(charset)
                    except UnicodeDecodeError:
                        body = payload.decode('latin1', errors='replace')
            
            # Format the email content
            content = f"Subject: {subject}\n"
            content += f"From: {from_addr}\n"
            content += f"To: {to_addr}\n"
            content += f"Date: {date}\n"
            content += f"\n{body}\n"
            
            return content
            
        except Exception as e:
            logger.error(f"Error extracting content from email {file_path}: {e}")
            # Fall back to reading as text
            return self._extract_text_content(file_path)
    
    def _extract_metadata(self, file_path: Path, file_type: str) -> Dict[str, Any]:
        """Extract metadata from file."""
        metadata = {
            "file_path": str(file_path),
            "file_name": file_path.name,
            "file_type": file_type,
            "extension": file_path.suffix.lower(),
            "size_bytes": file_path.stat().st_size,
            "created_at": datetime.fromtimestamp(file_path.stat().st_ctime).isoformat(),
            "modified_at": datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
        }
        
        # Extract additional metadata based on file type
        if file_type == "pdf" and "pdf" in self._parsers:
            try:
                import pypdf
                with open(file_path, 'rb') as f:
                    pdf = pypdf.PdfReader(f)
                    info = pdf.metadata
                    if info:
                        pdf_info = {}
                        for key, value in info.items():
                            if key.startswith('/'):
                                key = key[1:]
                            pdf_info[key] = str(value)
                        metadata["pdf_info"] = pdf_info
                    metadata["page_count"] = len(pdf.pages)
            except Exception as e:
                logger.error(f"Error extracting PDF metadata {file_path}: {e}")
        
        elif file_type == "docx" and "docx" in self._parsers:
            try:
                import docx
                doc = docx.Document(file_path)
                metadata["page_count"] = len(doc.sections)
                core_props = doc.core_properties
                metadata["docx_info"] = {
                    "title": core_props.title,
                    "author": core_props.author,
                    "created": core_props.created.isoformat() if core_props.created else None,
                    "modified": core_props.modified.isoformat() if core_props.modified else None,
                    "last_modified_by": core_props.last_modified_by,
                    "revision": core_props.revision,
                    "word_count": len(doc.paragraphs),
                }
            except Exception as e:
                logger.error(f"Error extracting DOCX metadata {file_path}: {e}")
        
        return metadata
    
    def _process_file(self, file_path: Path) -> Optional[Document]:
        """Process a file into a Document."""
        try:
            # Determine file type
            file_type = self._get_file_type(file_path)
            
            # Check if we have a parser for this file type
            if file_type not in self._parsers:
                logger.warning(f"No parser for file type: {file_type} ({file_path})")
                return None
            
            # Extract content
            content = self._parsers[file_type](file_path)
            if not content.strip():
                logger.warning(f"Empty content extracted from {file_path}")
                return None
            
            # Extract metadata
            metadata = self._extract_metadata(file_path, file_type) if self.extract_metadata else {
                "file_path": str(file_path),
                "file_name": file_path.name,
                "file_type": file_type
            }
            
            # Create source ID from file path
            file_path_str = str(file_path)
            source_id = file_path_str.replace('/', '_')
            source_id = source_id.replace('\\', '_')
            
            # Get creation and modification times
            created_at = datetime.fromtimestamp(file_path.stat().st_ctime)
            updated_at = datetime.fromtimestamp(file_path.stat().st_mtime)
            
            return Document(
                content=content,
                metadata=metadata,
                source_type=file_type,
                source_id=source_id,
                created_at=created_at,
                updated_at=updated_at
            )
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}")
            return None
    
    def _filter_files_by_query(self, files: List[Path], query: str) -> List[Path]:
        """Filter files based on a search query."""
        filtered_files = []
        
        for file_path in files:
            # Check filename
            if query.lower() in file_path.name.lower():
                filtered_files.append(file_path)
                continue
            
            # Check content (first 1000 bytes)
            try:
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    first_chunk = f.read(1000)
                    if query.lower() in first_chunk.lower():
                        filtered_files.append(file_path)
            except UnicodeDecodeError:
                # For binary files
                try:
                    with open(file_path, 'rb') as f:
                        first_chunk = f.read(1000).decode('utf-8', errors='replace')
                        if query.lower() in first_chunk.lower():
                            filtered_files.append(file_path)
                except Exception:
                    pass
            except Exception:
                pass
        
        return filtered_files
    
    async def get_documents(self, query: Optional[str] = None, 
                          limit: Optional[int] = None,
                          **kwargs) -> List[Document]:
        """
        Retrieve documents from files.
        
        Args:
            query: Optional query to filter files
            limit: Maximum number of documents to retrieve
            **kwargs: Additional parameters:
                - base_dirs: List of directories to scan
                - file_types: List of file types to include
                - patterns: List of glob patterns to include
                - recursive: Whether to scan recursively
        
        Returns:
            List of Document objects
        """
        if not await self.connect():
            return []
        
        try:
            # Update config with any runtime parameters
            if "base_dirs" in kwargs:
                self.base_dirs = [Path(p) for p in kwargs["base_dirs"]]
            if "file_types" in kwargs:
                self.extensions = kwargs["file_types"]
            if "patterns" in kwargs:
                self.include_patterns = kwargs["patterns"]
            if "recursive" in kwargs:
                self.recursive = kwargs["recursive"]
            
            limit = limit or self.max_files
            
            # Get all matching files
            all_files = self._get_all_files()
            
            # Filter files by query if provided
            if query:
                all_files = self._filter_files_by_query(all_files, query)
            
            # Limit the number of files
            all_files = all_files[:limit]
            
            # Process files into documents
            documents = []
            for file_path in all_files:
                document = self._process_file(file_path)
                if document:
                    documents.append(document)
                
                # Stop if we've reached the limit
                if len(documents) >= limit:
                    break
            
            return documents
            
        except Exception as e:
            logger.error(f"Error retrieving documents from files: {e}")
            return []
    
    async def stream_documents(self, query: Optional[str] = None, 
                             **kwargs) -> Generator[Document, None, None]:
        """
        Stream documents from files.
        
        Args:
            query: Optional query to filter files
            **kwargs: Additional parameters (same as get_documents)
            
        Yields:
            Document objects one at a time
        """
        documents = await self.get_documents(query, **kwargs)
        for doc in documents:
            yield doc
    
    async def get_document_by_id(self, doc_id: str) -> Optional[Document]:
        """
        Retrieve a specific document by ID.
        
        Args:
            doc_id: Document ID (based on file path)
            
        Returns:
            Document object if found, None otherwise
        """
        # Convert doc_id back to file path
        file_path = Path(doc_id.replace('_', '/'))
        
        # Check if file exists
        if not file_path.exists() or not file_path.is_file():
            # Try alternate conversion (for Windows paths)
            file_path = Path(doc_id.replace('_', '\\'))
            if not file_path.exists() or not file_path.is_file():
                return None
        
        # Process the file
        return self._process_file(file_path)